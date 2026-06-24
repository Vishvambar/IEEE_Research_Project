import docx2txt
from pptx import Presentation
import os

def extract_docx(filepath, outpath):
    text = docx2txt.process(filepath)
    with open(outpath, "w") as f:
        f.write(text)
    print(f"Extracted {filepath} to {outpath}")

def extract_pptx(filepath, outpath):
    prs = Presentation(filepath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    
    with open(outpath, "w") as f:
        f.write("\n".join(text_runs))
    print(f"Extracted {filepath} to {outpath}")

extract_docx("Research docs/Ready to Submit/Week_1_Healthcare_group_4.docx", "week1_report.md")
extract_docx("Research docs/Templates/Week_2_Template.docx", "week2_template.md")
extract_pptx("Research docs/Templates/IEEE Week 2 Poster Template.pptx", "week2_poster.md")
